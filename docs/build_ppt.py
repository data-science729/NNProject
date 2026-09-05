# -*- coding: utf-8 -*-
"""
综合课程设计报告第一部分 PPT 构建脚本 (逐步构建版)
当前状态：已完成第六页、第七页（研发困惑精简解答 QA 模块）
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_text_formatting(paragraph, text, font_name="Microsoft YaHei", font_size=Pt(14), color=RGBColor(45, 55, 72), bold=False, line_spacing=1.15):
    paragraph.text = text
    paragraph.font.name = font_name
    paragraph.font.size = font_size
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.line_spacing = line_spacing

def add_slide_header(slide, title_text):
    """为每一页内容页添加统一的精美页眉和蓝色横条装饰"""
    # 顶部深蓝色背景饰条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(13.33), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(26, 54, 93) # 深蓝色
    shape.line.fill.background()
    
    # 标题文字
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    apply_text_formatting(p, title_text, font_size=Pt(24), color=RGBColor(255, 255, 255), bold=True)

def main():
    prs = Presentation()
    # 设置宽屏 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 使用空白幻灯片布局
    blank_layout = prs.slide_layouts[6]
    
    # ==========================================
    # SLIDE 1: 封面页
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(26, 54, 93)
    bg_shape.line.fill.background()
    
    stripe = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.8))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(49, 151, 149)
    stripe.line.fill.background()
    
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(2.2))
    tf = title_box.text_frame
    p_title = tf.paragraphs[0]
    apply_text_formatting(p_title, "综合课程设计报告", font_size=Pt(44), color=RGBColor(255, 255, 255), bold=True)
    p_title.space_after = Pt(12)
    p_subtitle = tf.add_paragraph()
    apply_text_formatting(p_subtitle, "第一部分：数据预处理与中文文本分词系统设计", font_size=Pt(20), color=RGBColor(200, 230, 230))
    
    info_box = slide1.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(11.0), Inches(2.2))
    tf_info = info_box.text_frame
    p_class = tf_info.paragraphs[0]
    apply_text_formatting(p_class, "专业班级：2024级数据科学与大数据技术08班  |  指导老师：李筱光", font_size=Pt(15), color=RGBColor(220, 225, 230))
    p_class.space_after = Pt(8)
    p_members = tf_info.add_paragraph()
    apply_text_formatting(p_members, "小组成员：刘宇翔（组长）  王振  邸凯硕  王杨浩", font_size=Pt(15), color=RGBColor(220, 225, 230))
    
    # ==========================================
    # SLIDE 2: NLP 处理全流程与实例展示
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "一、 NLP 中文文本处理全生命周期")
    
    left_box = slide2.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p_l_header = tf_l.paragraphs[0]
    apply_text_formatting(p_l_header, "NLP 数据处理核心管道 (Pipeline)", font_size=Pt(18), color=RGBColor(26, 54, 93), bold=True)
    p_l_header.space_after = Pt(12)
    
    steps = [
        ("① Text (原始文本清洗与分词)", "去除 HTML 标签、杂质噪声，利用 Jieba 分词将长文本切分为单词列表，过滤停用词与低频词。"),
        ("② Dictionary (词典构建与剪枝)", "扫描整个语料库建立唯一的映射关系（str ↔ int）。通过剪枝（Pruning）剔除极罕见或极常见词，压缩特征维度，防范 OOM 异常。"),
        ("③ Corpus (数字化稀疏表示/词袋)", "将文档转化为词袋模型（Bag of Words）。统计每个单词 ID 在当前文档中出现的频数，表示为 [(词ID, 词频), ...] 的嵌套列表。"),
        ("④ Model (语义模型转换)", "克服词袋模型维度过高且忽略词间语义的缺点。利用数学模型（如 TF-IDF 算法）将“词频”转化为表示词汇核心重要度的“特征权重”。"),
        ("⑤ Vector Space (向量空间坐标)", "将文档转化到最终的连续向量空间中。表示为 [(词ID, 权重), ...]，用权重矩阵直接对接下游深度学习算法模型。")
    ]
    for i, (title, desc) in enumerate(steps):
        p_title = tf_l.add_paragraph()
        apply_text_formatting(p_title, title, font_size=Pt(13), color=RGBColor(49, 151, 149), bold=True)
        p_title.space_before = Pt(6)
        p_desc = tf_l.add_paragraph()
        apply_text_formatting(p_desc, desc, font_size=Pt(11), color=RGBColor(74, 85, 104))
        p_desc.space_before = Pt(2)
        
    right_box = slide2.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.5))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r_header = tf_r.paragraphs[0]
    apply_text_formatting(p_r_header, "笔记实例步骤图解 (Example Walkthrough)", font_size=Pt(18), color=RGBColor(26, 54, 93), bold=True)
    p_r_header.space_after = Pt(12)
    
    example_flow = [
        ("1. Text (输入两篇已分词新闻)", "Doc 1: ['大数据', '分析', '技术']\nDoc 2: ['自然语言', '处理', '技术']"),
        ("2. Dictionary (建立映射字典)", '{"分析": 0, "大数据": 1, "技术": 2, "处理": 3, "自然语言": 4}'),
        ("3. Corpus (数字化词袋表示)", "Doc 1: [(0, 1), (1, 1), (2, 1)]\nDoc 2: [(2, 1), (3, 1), (4, 1)]"),
        ("4. Model (应用 TF-IDF 算法转换)", "计算逆文档频率。因“技术”一词在两篇文档中都出现，不具备区分度，其 IDF 权重被降低过滤。"),
        ("5. Vector Space (得到向量空间最终投影)", "Doc 1: [(0, 0.7071), (1, 0.7071)]\nDoc 2: [(3, 0.7071), (4, 0.7071)]\n(词ID 2 的“技术”已被降权过滤，突出了大数据、分析、处理等独有特征)")
    ]
    for i, (stage, val) in enumerate(example_flow):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.1 + i * 1.0), Inches(5.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(247, 250, 252)
        card.line.color.rgb = RGBColor(226, 232, 240)
        card_text_box = slide2.shapes.add_textbox(Inches(7.1), Inches(2.1 + i * 1.0), Inches(5.5), Inches(0.85))
        tf_c = card_text_box.text_frame
        tf_c.word_wrap = True
        p_stage = tf_c.paragraphs[0]
        apply_text_formatting(p_stage, stage, font_size=Pt(11), color=RGBColor(26, 54, 93), bold=True)
        p_val = tf_c.add_paragraph()
        font_sz = Pt(9.5) if "\n" in val else Pt(10)
        apply_text_formatting(p_val, val, font_size=font_sz, color=RGBColor(45, 55, 72))

    # ==========================================
    # SLIDE 3: Jieba 分词模式与应用
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "二、 中文分词工具：Jieba分词机制与应用")
    
    left_box3 = slide3.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5))
    tf_l3 = left_box3.text_frame
    tf_l3.word_wrap = True
    p_l3_header = tf_l3.paragraphs[0]
    apply_text_formatting(p_l3_header, "Jieba 分词的三大核心模式", font_size=Pt(18), color=RGBColor(26, 54, 93), bold=True)
    p_l3_header.space_after = Pt(15)
    
    modes_info = [
        ("① 精确模式 (Accurate Mode)", "【原理】试图将句子最精确地切开，适合文本分析和 NLP 计算。\n【代码】jieba.cut(text, cut_all=False)"),
        ("② 全模式 (Full Mode)", "【原理】扫描句子中所有可以成词的词语。速度极快，但无法解决歧义。\n【代码】jieba.cut(text, cut_all=True)\n【特点】适用于信息检索，尽可能多地涵盖关键词，提高召回率。"),
        ("③ 搜索引擎模式 (Search Engine Mode)", "【原理】在精确模式的基础上，对长词再次切分。提高召回率。\n【代码】jieba.cut_for_search(text)\n【特点】适用于搜索引擎构建倒排索引。")
    ]
    for i, (title, desc) in enumerate(modes_info):
        p_t = tf_l3.add_paragraph()
        apply_text_formatting(p_t, title, font_size=Pt(14), color=RGBColor(49, 151, 149), bold=True)
        p_t.space_before = Pt(8)
        p_d = tf_l3.add_paragraph()
        apply_text_formatting(p_d, desc, font_size=Pt(11.5), color=RGBColor(74, 85, 104))
        p_d.space_before = Pt(3)

    right_box3 = slide3.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.5))
    tf_r3 = right_box3.text_frame
    tf_r3.word_wrap = True
    p_r3_header = tf_r3.paragraphs[0]
    apply_text_formatting(p_r3_header, "本项目中的分词应用细节", font_size=Pt(18), color=RGBColor(26, 54, 93), bold=True)
    p_r3_header.space_after = Pt(15)
    
    card3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.2), Inches(5.7), Inches(4.5))
    card3.fill.solid()
    card3.fill.fore_color.rgb = RGBColor(247, 250, 252)
    card3.line.color.rgb = RGBColor(226, 232, 240)
    
    card3_text_box = slide3.shapes.add_textbox(Inches(7.2), Inches(2.3), Inches(5.3), Inches(4.3))
    tf_c3 = card3_text_box.text_frame
    tf_c3.word_wrap = True
    p_c3_title = tf_c3.paragraphs[0]
    apply_text_formatting(p_c3_title, "★ 我们选择的模式：精确模式 (Accurate Mode)", font_size=Pt(14), color=RGBColor(26, 54, 93), bold=True)
    p_c3_title.space_after = Pt(12)
    p_c3_func = tf_c3.add_paragraph()
    apply_text_formatting(p_c3_func, "【调用函数】jieba.lcut(text)", font_size=Pt(13), color=RGBColor(49, 151, 149), bold=True)
    p_c3_func.space_after = Pt(8)
    p_c3_desc = tf_c3.add_paragraph()
    apply_text_formatting(p_c3_desc, "【设计目的与优势】\n"
                                     "1. 直接返回列表：jieba.lcut 函数会在分词后，将各个独立的词汇直接以 List (列表) 的形式返回。这是中文文本“Tokenize (词元化)”的最优实现。\n\n"
                                     "2. 避免语义噪音：与全模式相比，精确模式只对每个字切分一次，能够有效避免全模式带来的词汇交错与“冗余特征叠加”，保证了分类语义特征的纯净度。\n\n"
                                     "3. 完美对接词典：分词后的列表可直接被用于构建 word_vocab.json，以及计算每个词汇的一热编码索引。",
                          font_size=Pt(12), color=RGBColor(74, 85, 104))

    # ==========================================
    # SLIDE 4: 预处理成果总览与下游交付
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "三、 数据处理阶段成果总结与下游交付方案")
    
    left_box4 = slide4.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.5))
    tf_l4 = left_box4.text_frame
    tf_l4.word_wrap = True
    p_l4_header = tf_l4.paragraphs[0]
    apply_text_formatting(p_l4_header, "预处理核心代码与产出成果", font_size=Pt(18), color=RGBColor(26, 54, 93), bold=True)
    p_l4_header.space_after = Pt(12)
    
    p_code_t = tf_l4.add_paragraph()
    apply_text_formatting(p_code_t, "■ 团队编写的预处理核心代码：", font_size=Pt(14), color=RGBColor(49, 151, 149), bold=True)
    p_code_d = tf_l4.add_paragraph()
    apply_text_formatting(p_code_d, "1. src/data_utils.py：封装 TextPreprocessor 工具类，提供清洗、分词、截断填充和 PyTorch Dataset 数据加载接口。\n"
                                    "2. run_preprocessing.py：执行脚本，一键批量读取原始文本、调用清洗分词流并输出干净数据集。", font_size=Pt(11.5), color=RGBColor(74, 85, 104))
    p_code_d.space_after = Pt(10)
    
    p_files_t = tf_l4.add_paragraph()
    apply_text_formatting(p_files_t, "■ 批量执行后得到的四个核心文件：", font_size=Pt(14), color=RGBColor(49, 151, 149), bold=True)
    
    files_list = [
        ("cnews.train.clean.txt", "包含 50,000 条经过分词清洗的训练集文件"),
        ("cnews.val.clean.txt", "包含 5,000 条分词清洗的验证集文件"),
        ("cnews.test.clean.txt", "包含 10,000 条分词清洗的测试集文件"),
        ("word_vocab.json", "包含 10,000 个高频词一热映射的全局索引字典")
    ]
    for fname, fdesc in files_list:
        p_item = tf_l4.add_paragraph()
        apply_text_formatting(p_item, f"• {fname}：{fdesc}", font_size=Pt(11), color=RGBColor(74, 85, 104))
        p_item.space_before = Pt(2)

    right_box4 = slide4.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.5))
    tf_r4 = right_box4.text_frame
    tf_r4.word_wrap = True
    p_r4_header = tf_r4.paragraphs[0]
    apply_text_formatting(p_r4_header, "解耦交付：让后级组员工作高效启动", font_size=Pt(18), color=RGBColor(26, 54, 93), bold=True)
    p_r4_header.space_after = Pt(12)
    
    card4_1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.9), Inches(2.1))
    card4_1.fill.solid()
    card4_1.fill.fore_color.rgb = RGBColor(240, 244, 248)
    card4_1.line.color.rgb = RGBColor(200, 214, 224)
    
    tb4_1 = slide4.shapes.add_textbox(Inches(6.9), Inches(2.25), Inches(5.7), Inches(2.0))
    tf4_1 = tb4_1.text_frame
    tf4_1.word_wrap = True
    p1_t = tf4_1.paragraphs[0]
    apply_text_formatting(p1_t, "📂 交付物 AtoB ➔ 第二组（同学 B · 词向量训练）", font_size=Pt(12), color=RGBColor(26, 54, 93), bold=True)
    p1_d = tf4_1.add_paragraph()
    apply_text_formatting(p1_d, "【包含文件】cnews.train.clean.txt + word_vocab.json\n"
                                "【作用与目的】\n"
                                "1. 提供清洗分词好的干净语料，直接输入训练 Word2Vec CBOW 模型。\n"
                                "2. 提供全局词表，用于将词向量按 ID 对齐，生成 10000x100 的 Embedding 权重矩阵。", font_size=Pt(10.5), color=RGBColor(74, 85, 104))
    p1_d.space_before = Pt(4)
    
    card4_2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.5), Inches(5.9), Inches(2.4))
    card4_2.fill.solid()
    card4_2.fill.fore_color.rgb = RGBColor(230, 242, 242)
    card4_2.line.color.rgb = RGBColor(178, 204, 204)
    
    tb4_2 = slide4.shapes.add_textbox(Inches(6.9), Inches(4.55), Inches(5.7), Inches(2.3))
    tf4_2 = tb4_2.text_frame
    tf4_2.word_wrap = True
    p2_t = tf4_2.paragraphs[0]
    apply_text_formatting(p2_t, "📂 交付物 AtoC ➔ 第三组（同学 C · 神经网络分类）", font_size=Pt(12), color=RGBColor(26, 54, 93), bold=True)
    p2_d = tf4_2.add_paragraph()
    apply_text_formatting(p2_d, "【包含文件】src/data_utils.py + word_vocab.json + 所有三个 .clean.txt 文件\n"
                                "【作用与目的】\n"
                                "1. 提供统一的数据集封装（THUCNewsDataset 类），秒级加载，兼容批处理。\n"
                                "2. 提供干净的划分好的训练集和验证集，用于模型拟合，让 C 组无需关注底层清洗工作，直接启动模型开发。", font_size=Pt(10.5), color=RGBColor(74, 85, 104))
    p2_d.space_before = Pt(4)

    # ==========================================
    # SLIDE 5: TextPreprocessor 算法设计与伪代码
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "四、 预处理器 TextPreprocessor 算法伪代码与设计")
    
    code_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.3))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
    code_bg.line.fill.background()
    
    code_box = slide5.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(6.3), Inches(5.1))
    tf_code = code_box.text_frame
    tf_code.word_wrap = True
    
    p_code = tf_code.paragraphs[0]
    apply_text_formatting(p_code, "class TextPreprocessor (核心文本预处理器) 伪代码", font_size=Pt(12), color=RGBColor(78, 201, 176), bold=True)
    
    code_lines = [
        "  def clean_text(raw_text):",
        "      1. 正则过滤HTML标签: re.sub('<[^>]+>', '', raw_text)",
        "      2. 保留常用汉字(u4e00-u9fa5)、英文、数字及中文句读标点",
        "         过滤其余无意义噪音字符 -> 返回 clean_text",
        "",
        "  def segment_text(clean_text):",
        "      1. 精确分词: words = jieba.lcut(clean_text)",
        "      2. 过滤 stopwords.txt (停用词表) 中的无语义词",
        "         -> 返回干净词列表 words",
        "",
        "  def build_vocab(corpus, vocab_size=10000):",
        "      1. 统计整个训练语料的词频，过滤总频数 < 5 的长尾词",
        "      2. 截取最高频的前 9998 个词，注入首尾预设的占位符：",
        "         word2id = {'<PAD>': 0, '<UNK>': 1, ... 高频词映射}",
        "         -> 保存为 word_vocab.json 对照表",
        "",
        "  def text_to_ids(text):",
        "      1. 查表编码: ids = [word2id.get(w, 1) for w in text]",
        "      2. 填充与截断: 若长度 < 888 用 0(<PAD>) 填充补齐;",
        "                    若长度 > 888 则截断 -> 返回 888 维 ID 列表"
    ]
    for cl in code_lines:
        p_l = tf_code.add_paragraph()
        color = RGBColor(220, 220, 170) if "def " in cl else RGBColor(214, 157, 133) if "1." in cl or "2." in cl else RGBColor(200, 200, 200)
        apply_text_formatting(p_l, cl, font_size=Pt(9.5), font_name="Consolas", color=color)
        p_l.space_after = Pt(2)

    right_box5 = slide5.shapes.add_textbox(Inches(7.3), Inches(1.5), Inches(5.4), Inches(5.3))
    tf_r5 = right_box5.text_frame
    tf_r5.word_wrap = True
    p_r5_header = tf_r5.paragraphs[0]
    apply_text_formatting(p_r5_header, "算法设计理念与核心逻辑解读", font_size=Pt(18), color=RGBColor(26, 54, 93), bold=True)
    p_r5_header.space_after = Pt(15)
    
    interpretations = [
        ("★ 模块化与流水线封装 (Pipeline)", "将数据清洗、分词、词汇编码（一热对照）及定长对齐逻辑完全解耦并集成在一个类中，为下游批量处理提供了极佳 of 工程通用性。"),
        ("★ 鲁棒的 OOV 未登录词处理机制", "在 text_to_ids 方法中，使用字典的 .get(w, 1) 进行安全查表。如果遇到生僻词，安全地fallback至 ID 1 (即 <UNK>)，有效规避了运行时 KeyError 崩溃故障。"),
        ("★ 数据并行与形状一致性要求", "通过定长对齐逻辑，将任意长度文本标准化映射为 888 维一维整型向量。这是下游 PyTorch 进行高效 DataLoader 批数据生成和 GPU 并行矩阵计算的前提。")
    ]
    for i, (title, text) in enumerate(interpretations):
        p_t = tf_r5.add_paragraph()
        apply_text_formatting(p_t, title, font_size=Pt(13), color=RGBColor(49, 151, 149), bold=True)
        p_t.space_before = Pt(8)
        p_d = tf_r5.add_paragraph()
        apply_text_formatting(p_d, text, font_size=Pt(11.5), color=RGBColor(74, 85, 104))
        p_d.space_before = Pt(3)

    # ==========================================
    # SLIDE 6: 研发过程困惑解答 QA (一) - 标点语义与特殊字符 (3 张垂直卡片排版)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "五、 研发核心疑难解答 (Q&A) — 标点语义与特殊字符")
    
    # 3张垂直卡片
    qa_list6 = [
        ("Q1: 为什么清洗阶段保留标点，但停用词阶段又过滤它们？这是否矛盾？",
         "【解答】不矛盾。这属于“分段式语境保留”机制：\n"
         "1. 保留标点（，。！？：）是为了让分词器能够以句读为边界，获取更准确的“局部上下文词组关系”，这直接有利于 Word2Vec 捕获句内关联；\n"
         "2. 在分词后，将它们在停用词阶段剔除，是因为标点在矩阵层面没有分类指向性，过滤它们可以排除噪声、极致精简特征编码向量的有效维度。"),
         
        ("Q2: 词表首尾手动注入的 <PAD> 和 <UNK> 占位符分别有什么作用？",
         "【解答】这两个符号是确保下游计算和推理能够顺利进行的底层护航者：\n"
         "1. <PAD> (ID 0 - 填充符)：解决新闻文本长短不一的问题。通过在短新闻尾部拼装 0，将批数据统一规格化为 [batch_size, 888] 形状以符合 GPU 矩阵乘法的对齐要求；\n"
         "2. <UNK> (ID 1 - 未知符)：兜底未知词。测试集或最终应用中出现的生僻词/错字会统一归为 ID 1，防止字典查询触发 KeyError 崩溃。"),
         
        ("Q3: 停用词 (Stopwords) 的主要内容和核心过滤目的是什么？",
         "【解答】主要内容与过滤机制如下：\n"
         "1. 停用词通常包含无具体实际分类属性的中文结构助词（如的、地、得）、介词（在、于）、人称代词和强连接连词；\n"
         "2. 过滤停用词的最核心目的是降低数据的“噪音背景”。如果不加过滤，高频的无语义词会霸占前 10,000 的词表空间，导致有区分度的实词特征（如“涨停”、“扣篮”）被稀释，严重干扰网络的注意力。")
    ]
    
    for i, (q_text, a_text) in enumerate(qa_list6):
        # 绘制背景卡片
        card = slide6.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), Inches(1.5 + i * 1.85), Inches(12.13), Inches(1.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        # 写入文本
        tb = slide6.shapes.add_textbox(Inches(0.7), Inches(1.5 + i * 1.85), Inches(11.93), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pq = tf.paragraphs[0]
        apply_text_formatting(pq, q_text, font_size=Pt(12.5), color=RGBColor(26, 54, 93), bold=True)
        pq.space_after = Pt(4)
        
        pa = tf.add_paragraph()
        apply_text_formatting(pa, a_text, font_size=Pt(10.5), color=RGBColor(74, 85, 104))

    # ==========================================
    # SLIDE 7: 研发过程困惑解答 QA (二) - 超参数配置机制 (左右双栏卡片排版)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "六、 研发核心疑难解答 (Q&A) — 核心超参数配置原理")
    
    # 左栏卡片：min_freq
    card_l7 = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.1)
    )
    card_l7.fill.solid()
    card_l7.fill.fore_color.rgb = RGBColor(240, 244, 248) # 浅灰蓝色
    card_l7.line.color.rgb = RGBColor(200, 214, 224)
    
    tb_l7 = slide7.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(5.6), Inches(4.9))
    tf_l7 = tb_l7.text_frame
    tf_l7.word_wrap = True
    
    p_l7_q = tf_l7.paragraphs[0]
    apply_text_formatting(p_l7_q, "Q4: 为什么过滤阈值 min_freq 设为 5？\n如果直接保留全部词有什么后果？", font_size=Pt(14), color=RGBColor(26, 54, 93), bold=True)
    p_l7_q.space_after = Pt(12)
    
    p_l7_a = tf_l7.add_paragraph()
    apply_text_formatting(p_l7_a, "【科学依据与危害剖析】\n\n"
                                  "1. 过滤罕见噪声：在 5 万篇训练集中出现小于 5 次的词（如错字、特殊ID）通常没有统计分类价值。将其过滤能显著提高词表的信噪比；\n\n"
                                  "2. 规避过拟合：若不做限制而保留只出现一次的错别字，神经网络会产生硬背效应（即把噪音错字绑定为类别特征），严重破坏泛化精度；\n\n"
                                  "3. 减小特征维度：若不进行频数过滤，词表大小将从 10,000 跃升至 364,311 维。这会导致 Embedding 词向量权重参数膨胀 36 倍以上，直接引发显存溢出（OOM）崩溃，严重拖慢训练收敛。",
                          font_size=Pt(12), color=RGBColor(74, 85, 104))

    # 右栏卡片：max_len
    card_r7 = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.1)
    )
    card_r7.fill.solid()
    card_r7.fill.fore_color.rgb = RGBColor(230, 242, 242) # 浅灰绿色
    card_r7.line.color.rgb = RGBColor(178, 204, 204)
    
    tb_r7 = slide7.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.6), Inches(4.9))
    tf_r7 = tb_r7.text_frame
    tf_r7.word_wrap = True
    
    p_r7_q = tf_r7.paragraphs[0]
    apply_text_formatting(p_r7_q, "Q5: 文本序列对齐长度 max_len 选择 888 的依据是什么？过长或过短有何不利？", font_size=Pt(14), color=RGBColor(26, 54, 93), bold=True)
    p_r7_q.space_after = Pt(12)
    
    p_r7_a = tf_r7.add_paragraph()
    apply_text_formatting(p_r7_a, "【统计学与硬件性能分析】\n\n"
                                  "1. 统计学 90 分位数：在对 THUCNews 数据集进行分词长度统计后发现，90% 以上的新闻正文长度均包含在 888 个词以内。888 是保证信息无损读入的最优分界点；\n\n"
                                  "2. 避免信息丢失：若强行设小（如限制为 60 词），模型在提取特征时将只看导语而强行丢失文章中后段包含的大量关键分类词汇，导致准确率骤降；\n\n"
                                  "3. 控制计算时序损耗：LSTM 等模型的计算开销和时序长度成正比。若设得太大（如 6000 词），较短的文本必须追加大量的 0，导致显卡 GPU 浪费大量算力在无效计算上，拖慢计算速度数倍。",
                          font_size=Pt(12), color=RGBColor(74, 85, 104))

    # 保存 PPT
    import time
    save_path = r"c:\Users\asus\PycharmProjects\NNProject\NN\综合课程设计报告第一部分.pptx"
    try:
        prs.save(save_path)
        print(f"[PPT] 第六页和第七页QA已成功追加并保存至: {save_path}")
    except PermissionError:
        backup_path = r"c:\Users\asus\PycharmProjects\NNProject\NN\综合课程设计报告第一部分_新.pptx"
        try:
            prs.save(backup_path)
            print(f"[PPT] 警告：主文件已被占用，已另存为: {backup_path}")
        except PermissionError:
            timestamp_path = f"c:\\Users\\asus\\PycharmProjects\\NNProject\\NN\\综合课程设计报告第一部分_{int(time.time())}.pptx"
            prs.save(timestamp_path)
            print(f"[PPT] 警告：常规PPT文件均被打开占用，已自动输出到独立版本文件: {timestamp_path}")

if __name__ == "__main__":
    main()
